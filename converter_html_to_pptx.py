#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Conversor HTML para PPTX - Apresentação Arboviroses
Converte a apresentação HTML para formato PowerPoint compatível com Canva
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import html

def hex_to_rgb(hex_color):
    """Converte cor hexadecimal para RGB"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def extract_slides_from_html(html_file):
    """Extrai conteúdo dos slides do arquivo HTML"""
    with open(html_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Encontrar todos os slides
    slide_pattern = r'<div class="slide">(.*?)</div>\s*(?=<div class="slide">|$)'
    slides = re.findall(slide_pattern, content, re.DOTALL)

    slides_data = []

    for i, slide_content in enumerate(slides):
        # Extrair número do slide
        slide_num_match = re.search(r'Slide (\d+)/42', slide_content)
        slide_num = slide_num_match.group(1) if slide_num_match else str(i+1)

        # Extrair título principal (h1 ou h2)
        title_match = re.search(r'<h[12][^>]*>(.*?)</h[12]>', slide_content, re.DOTALL)
        title = html.unescape(re.sub(r'<[^>]+>', '', title_match.group(1))) if title_match else f"Slide {slide_num}"

        # Extrair cabeçalho do aluno se existir
        aluno_match = re.search(r'<div class="aluno-header"[^>]*>(.*?)</div>', slide_content, re.DOTALL)
        aluno_info = html.unescape(re.sub(r'<[^>]+>', '', aluno_match.group(1))) if aluno_match else ""

        # Extrair conteúdo de cards
        card_pattern = r'<div class="card"[^>]*>(.*?)</div>'
        cards = re.findall(card_pattern, slide_content, re.DOTALL)

        card_contents = []
        for card in cards:
            # Extrair h3 do card
            h3_match = re.search(r'<h3[^>]*>(.*?)</h3>', card, re.DOTALL)
            card_title = html.unescape(re.sub(r'<[^>]+>', '', h3_match.group(1))) if h3_match else ""

            # Extrair lista ou conteúdo
            li_pattern = r'<li[^>]*>(.*?)</li>'
            items = re.findall(li_pattern, card, re.DOTALL)
            card_items = [html.unescape(re.sub(r'<[^>]+>', '', item)) for item in items]

            # Se não há lista, pegar divs com conteúdo
            if not card_items:
                div_pattern = r'<div[^>]*style="text-align: center[^>]*>(.*?)</div>'
                divs = re.findall(div_pattern, card, re.DOTALL)
                card_items = [html.unescape(re.sub(r'<[^>]+>', '', div)) for div in divs if div.strip()]

            if card_title or card_items:
                card_contents.append({
                    'title': card_title,
                    'items': card_items
                })

        # Extrair conteúdo de destaque (danger, highlight, stat-box)
        highlight_pattern = r'<div class="(?:danger|highlight|stat-box|danger-visual)"[^>]*>(.*?)</div>'
        highlights = re.findall(highlight_pattern, slide_content, re.DOTALL)
        highlight_texts = []
        for highlight in highlights:
            clean_text = html.unescape(re.sub(r'<[^>]+>', ' ', highlight))
            clean_text = re.sub(r'\s+', ' ', clean_text).strip()
            if clean_text:
                highlight_texts.append(clean_text)

        # Extrair referências
        ref_match = re.search(r'<div class="reference"[^>]*>(.*?)</div>', slide_content, re.DOTALL)
        reference = html.unescape(re.sub(r'<[^>]+>', '', ref_match.group(1))) if ref_match else ""

        slides_data.append({
            'number': slide_num,
            'title': title,
            'aluno_info': aluno_info,
            'cards': card_contents,
            'highlights': highlight_texts,
            'reference': reference
        })

    return slides_data

def create_pptx_from_slides(slides_data, output_file):
    """Cria arquivo PPTX a partir dos dados dos slides"""

    # Criar apresentação
    prs = Presentation()

    # Definir layout padrão
    slide_layout = prs.slide_layouts[6]  # Layout em branco

    for slide_data in slides_data:
        slide = prs.slides.add_slide(slide_layout)

        # Título principal
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = slide_data['title']
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(44, 62, 80)  # Azul escuro
        title_p.alignment = PP_ALIGN.CENTER

        y_position = Inches(1.2)

        # Informações do aluno
        if slide_data['aluno_info']:
            aluno_box = slide.shapes.add_textbox(Inches(0.5), y_position, Inches(9), Inches(0.5))
            aluno_frame = aluno_box.text_frame
            aluno_frame.text = slide_data['aluno_info']
            aluno_p = aluno_frame.paragraphs[0]
            aluno_p.font.size = Pt(14)
            aluno_p.font.bold = True
            aluno_p.font.color.rgb = RGBColor(231, 76, 60)  # Vermelho
            aluno_p.alignment = PP_ALIGN.CENTER
            y_position += Inches(0.7)

        # Cards em grid
        if slide_data['cards']:
            cards_per_row = min(3, len(slide_data['cards']))
            card_width = Inches(9) / cards_per_row - Inches(0.1)

            for i, card in enumerate(slide_data['cards']):
                row = i // cards_per_row
                col = i % cards_per_row

                x_pos = Inches(0.5) + col * (card_width + Inches(0.1))
                y_pos = y_position + row * Inches(2.2)

                # Caixa do card
                card_box = slide.shapes.add_textbox(x_pos, y_pos, card_width, Inches(2))
                card_frame = card_box.text_frame
                card_frame.word_wrap = True

                # Título do card
                if card['title']:
                    title_p = card_frame.paragraphs[0]
                    title_p.text = card['title']
                    title_p.font.size = Pt(14)
                    title_p.font.bold = True
                    title_p.font.color.rgb = RGBColor(52, 152, 219)  # Azul

                # Items do card
                for item in card['items']:
                    if item.strip():
                        p = card_frame.add_paragraph()
                        p.text = f"• {item}"
                        p.font.size = Pt(11)
                        p.font.color.rgb = RGBColor(44, 62, 80)  # Azul escuro
                        p.space_after = Pt(6)

        # Destaques
        if slide_data['highlights']:
            highlight_y = y_position + Inches(3) if slide_data['cards'] else y_position

            for highlight in slide_data['highlights']:
                if highlight.strip():
                    highlight_box = slide.shapes.add_textbox(Inches(1), highlight_y, Inches(8), Inches(0.8))
                    highlight_frame = highlight_box.text_frame
                    highlight_frame.text = highlight
                    highlight_p = highlight_frame.paragraphs[0]
                    highlight_p.font.size = Pt(12)
                    highlight_p.font.bold = True
                    highlight_p.font.color.rgb = RGBColor(231, 76, 60)  # Vermelho
                    highlight_p.alignment = PP_ALIGN.CENTER
                    highlight_y += Inches(0.9)

        # Referência no rodapé
        if slide_data['reference']:
            ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
            ref_frame = ref_box.text_frame
            ref_frame.text = slide_data['reference']
            ref_p = ref_frame.paragraphs[0]
            ref_p.font.size = Pt(9)
            ref_p.font.color.rgb = RGBColor(127, 140, 141)  # Cinza
            ref_p.alignment = PP_ALIGN.LEFT

        # Número do slide
        num_box = slide.shapes.add_textbox(Inches(8.5), Inches(0.1), Inches(1), Inches(0.3))
        num_frame = num_box.text_frame
        num_frame.text = f"Slide {slide_data['number']}/42"
        num_p = num_frame.paragraphs[0]
        num_p.font.size = Pt(10)
        num_p.font.color.rgb = RGBColor(255, 255, 255)  # Branco
        num_p.alignment = PP_ALIGN.RIGHT

        # Fundo colorido para número do slide
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(8.4), Inches(0.05), Inches(1.1), Inches(0.4)
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(52, 152, 219)  # Azul

    # Salvar apresentação
    prs.save(output_file)
    print(f"✅ Apresentação salva como: {output_file}")

def main():
    """Função principal"""
    html_file = "apresentação_v3.html"
    output_file = "apresentacao_arboviroses.pptx"

    print("🔄 Convertendo HTML para PPTX...")
    print(f"📁 Arquivo de entrada: {html_file}")
    print(f"📁 Arquivo de saída: {output_file}")

    try:
        # Extrair dados dos slides
        slides_data = extract_slides_from_html(html_file)
        print(f"📊 Encontrados {len(slides_data)} slides")

        # Criar PPTX
        create_pptx_from_slides(slides_data, output_file)

        print("✅ Conversão concluída com sucesso!")
        print(f"🎯 Agora você pode importar {output_file} no Canva")

    except Exception as e:
        print(f"❌ Erro durante a conversão: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    main()